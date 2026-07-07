"""Excel and PowerPoint exports of engine results + assessment.

Deterministic like report.py: every number comes from the results/assessment
objects. Excel ships every results table on its own sheet (SmartPLS-style
workbook); PowerPoint is a compact findings deck for supervision meetings.
"""
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

# results keys exported as sheets, in workbook order
RESULT_SHEETS = [
    ("loadings", "Loadings"),
    ("weights", "Outer weights"),
    ("reliability", "Reliability"),
    ("htmt", "HTMT"),
    ("fornell_larcker", "Fornell-Larcker"),
    ("cross_loadings", "Cross loadings"),
    ("paths_and_r2", "Paths and R2"),
    ("f_square", "f2"),
    ("total_effects", "Total effects"),
    ("total_indirect", "Total indirect"),
    ("specific_indirect", "Specific indirect"),
    ("blindfolding", "Blindfolding Q2"),
    ("prediction", "PLSpredict"),
    ("boot_paths", "Bootstrap paths"),
    ("boot_loadings", "Bootstrap loadings"),
    ("boot_weights", "Bootstrap weights"),
    ("boot_htmt", "Bootstrap HTMT"),
]


def _records_df(records) -> pd.DataFrame | None:
    if not isinstance(records, list) or not records:
        return None
    df = pd.DataFrame.from_records(records)
    return df.drop(columns=[c for c in df.columns if c.startswith("_")], errors="ignore")


def write_xlsx(path, results: dict, assessment: dict, mga: dict | None = None):
    """Full results workbook; one sheet per table."""
    with pd.ExcelWriter(path, engine="openpyxl") as xl:
        scalars = {"SRMR": results.get("srmr"), "NFI": results.get("nfi"),
                   "RMS_theta": results.get("rms_theta")}
        pd.DataFrame([{"metric": k, "value": v} for k, v in scalars.items()
                      if v is not None]).to_excel(xl, sheet_name="Model fit", index=False)
        for key, sheet in RESULT_SHEETS:
            df = _records_df(results.get(key))
            if df is not None:
                df.to_excel(xl, sheet_name=sheet, index=False)
        ipma = results.get("ipma") or {}
        for key, sheet in (("performance", "IPMA performance"),
                           ("total_effects_unstd", "IPMA importance")):
            df = _records_df(ipma.get(key))
            if df is not None:
                df.to_excel(xl, sheet_name=sheet, index=False)
        for key, sheet in (("measurement_model", "Assessment measurement"),
                           ("structural_model", "Assessment structural"),
                           ("hypotheses", "Hypotheses"),
                           ("mediation", "Mediation")):
            df = _records_df(assessment.get(key))
            if df is not None:
                df.to_excel(xl, sheet_name=sheet, index=False)
        if mga:
            for key, sheet in (("step2", "MICOM step 2"), ("step3", "MICOM step 3")):
                df = _records_df((mga.get("micom") or {}).get(key))
                if df is not None:
                    df.to_excel(xl, sheet_name=sheet, index=False)
            df = _records_df(mga.get("paths"))
            if df is not None:
                df.to_excel(xl, sheet_name="MGA paths", index=False)


def _slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = title
    return slide


def _pptx_table(slide, headers, rows, top=1.5):
    n_rows, n_cols = len(rows) + 1, len(headers)
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(top),
                                   Inches(9), Inches(0.35 * n_rows))
    table = shape.table
    for j, h in enumerate(headers):
        table.cell(0, j).text = h
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            cell = table.cell(i, j)
            cell.text = "" if v is None else str(v)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(12)
    return table


def _fmt(v):
    return f"{v:.3f}" if isinstance(v, (int, float)) else (v or "")


def write_pptx(path, dataset_meta: dict, request: dict, results: dict,
               assessment: dict, mga: dict | None = None):
    """Compact findings deck: model, verdicts, hypotheses, mediation, MGA."""
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "PLS-SEM Analysis — Key Findings"
    meta = results["meta"]
    slide.placeholders[1].text = (
        f"Dataset: {dataset_meta.get('filename')} (n = {meta['n']})\n"
        f"Engine: {meta['engine']} · {meta['nboot']:,} bootstrap resamples"
    )

    summ = assessment["summary"]
    slide = _slide(prs, "Assessment summary")
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
    tf = body.text_frame
    tf.text = (f"{summ['hypotheses_supported']} of {summ['hypotheses_total']} "
               f"hypothesized paths supported")
    for line in (
        f"Quality checks: {summ['pass']} pass · {summ['review']} review · {summ['fail']} fail",
        *(f"{s['metric']} = {_fmt(s['value'])} — {s['verdict'].upper()}"
          for s in assessment["structural_model"] if s["family"] == "model_fit"),
    ):
        tf.add_paragraph().text = line

    slide = _slide(prs, "Hypotheses")
    _pptx_table(slide, ["#", "Path", "β", "t", "Verdict"],
                [[h["hypothesis"], h["path"], _fmt(h["estimate"]),
                  _fmt(h["t_value"]), h["verdict"]]
                 for h in assessment["hypotheses"][:12]])

    mediation = assessment.get("mediation") or []
    if mediation:
        slide = _slide(prs, "Mediation — specific indirect effects")
        _pptx_table(slide, ["Indirect path", "β", "Classification"],
                    [[m["path"], _fmt(m["indirect_effect"]), m["classification"]]
                     for m in mediation[:12]])

    ipma = results.get("ipma") or {}
    perf = {r["row"]: r["performance"] for r in ipma.get("performance") or []}
    te = {r["row"]: r for r in ipma.get("total_effects_unstd") or []}
    outgoing = {p["from_construct"] for p in request["paths"]}
    endogenous = {p["to_construct"] for p in request["paths"]}
    targets = [c for c in perf if c in endogenous and c not in outgoing]
    if targets:
        target = targets[0]
        rows = sorted(
            ([p, _fmt(te.get(p, {}).get(target)), f"{perf[p]:.1f}"]
             for p in perf if p != target and isinstance(te.get(p, {}).get(target), (int, float))
             and abs(te[p][target]) > 1e-9),
            key=lambda r: -float(r[1]))
        if rows:
            slide = _slide(prs, f"IPMA — priorities for {target}")
            _pptx_table(slide, ["Construct", "Importance", "Performance (0–100)"], rows)

    if mga:
        m, micom = mga["meta"], mga["micom"]
        slide = _slide(prs, "Multi-group analysis")
        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.2))
        body.text_frame.text = (
            f"{m['group_variable']}: {m['value_a']} (n={m['n_a']}) vs "
            f"{m['value_b']} (n={m['n_b']}) — MICOM: {micom['invariance']} invariance")
        rows = [[p["path"], _fmt(p["estimate_a"]), _fmt(p["estimate_b"]),
                 _fmt(p["p_value"]), p["verdict"]] for p in mga["paths"][:10]]
        _pptx_table(slide, ["Path", "β (A)", "β (B)", "p", "Verdict"], rows, top=2.3)

    prs.save(path)
